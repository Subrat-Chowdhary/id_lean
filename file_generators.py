"""
File Generators for creating PPT and PDF outputs
"""

import os
from typing import List, Dict, Any, Optional
from pathlib import Path
import json

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# PDF generation
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors

# For charts and visualizations
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.backends.backend_pdf import PdfPages


def setup_matplotlib_for_plotting():
    """
    Setup matplotlib and seaborn for plotting with proper configuration.
    Call this function before creating any plots to ensure proper rendering.
    """
    import warnings
    
    # Ensure warnings are printed
    warnings.filterwarnings('default')  # Show all warnings

    # Configure matplotlib for non-interactive mode
    plt.switch_backend("Agg")

    # Set chart style
    plt.style.use("seaborn-v0_8")
    sns.set_palette("husl")

    # Configure platform-appropriate fonts for cross-platform compatibility
    # Must be set after style.use, otherwise will be overridden by style configuration
    plt.rcParams["font.sans-serif"] = ["Noto Sans CJK SC", "WenQuanYi Zen Hei", "PingFang SC", "Arial Unicode MS", "Hiragino Sans GB"]
    plt.rcParams["axes.unicode_minus"] = False


class PowerPointGenerator:
    """Generate PowerPoint presentations from content data"""
    
    def __init__(self, template_path: Optional[str] = None):
        self.template_path = template_path
        
        # Color scheme
        self.colors = {
            'primary': RGBColor(54, 96, 146),      # Professional blue
            'secondary': RGBColor(79, 129, 189),   # Lighter blue
            'accent': RGBColor(192, 80, 77),       # Red accent
            'text': RGBColor(64, 64, 64),          # Dark gray
            'light_gray': RGBColor(242, 242, 242)  # Light gray
        }
    
    def create_presentation(self, 
                          slides_data: List[Dict[str, Any]], 
                          output_path: str,
                          theme: str = 'professional') -> str:
        """
        Create PowerPoint presentation from slides data
        
        Args:
            slides_data: List of slide dictionaries
            output_path: Path to save the presentation
            theme: Visual theme ('professional', 'modern', 'minimal')
            
        Returns:
            Path to created presentation
        """
        # Create presentation object
        if self.template_path and os.path.exists(self.template_path):
            prs = Presentation(self.template_path)
        else:
            prs = Presentation()
        
        # Remove default slides
        while len(prs.slides) > 0:
            slide_id = prs.slides._slide_lst[0]
            prs.part.drop_rel(slide_id.rId)
            del prs.slides._slide_lst[0]
        
        # Create slides
        for slide_data in slides_data:
            self._create_slide(prs, slide_data, theme)
        
        # Save presentation
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        
        return output_path
    
    def _create_slide(self, prs: Presentation, slide_data: Dict[str, Any], theme: str):
        """Create individual slide based on type"""
        
        slide_type = slide_data.get('type', 'content')
        
        if slide_type == 'title':
            self._create_title_slide(prs, slide_data, theme)
        elif slide_type == 'bullet_points':
            self._create_bullet_slide(prs, slide_data, theme)
        elif slide_type == 'content':
            self._create_content_slide(prs, slide_data, theme)
        elif slide_type == 'activity':
            self._create_activity_slide(prs, slide_data, theme)
        elif slide_type == 'assessment':
            self._create_assessment_slide(prs, slide_data, theme)
        elif slide_type == 'conclusion':
            self._create_conclusion_slide(prs, slide_data, theme)
        else:
            self._create_generic_slide(prs, slide_data, theme)
    
    def _create_title_slide(self, prs: Presentation, slide_data: Dict[str, Any], theme: str):
        """Create title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = slide_data['content']['main_title']
        self._format_title_text(title, theme)
        
        # Subtitle
        if slide.placeholders[1]:
            subtitle = slide.placeholders[1]
            subtitle_text = f"{slide_data['content']['subtitle']}\n{slide_data['content']['duration']}\n{slide_data['content']['date']}"
            subtitle.text = subtitle_text
            self._format_subtitle_text(subtitle, theme)
    
    def _create_bullet_slide(self, prs: Presentation, slide_data: Dict[str, Any], theme: str):
        """Create bullet points slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = slide_data['title']
        self._format_title_text(title, theme)
        
        # Content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        
        # Add intro text if available
        if 'title' in slide_data['content']:
            p = text_frame.paragraphs[0]
            p.text = slide_data['content']['title']
            p.level = 0
            self._format_paragraph(p, theme, is_intro=True)
            
            # Add bullet points
            for bullet in slide_data['content']['bullet_points']:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 1
                self._format_paragraph(p, theme)
        else:
            # Direct bullet points
            p = text_frame.paragraphs[0]
            p.text = slide_data['content']['bullet_points'][0]
            p.level = 0
            self._format_paragraph(p, theme)
            
            for bullet in slide_data['content']['bullet_points'][1:]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                self._format_paragraph(p, theme)
    
    def _create_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], theme: str):
        """Create content slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = slide_data['title']
        self._format_title_text(title, theme)
        
        # Content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        
        # Main concept
        p = text_frame.paragraphs[0]
        p.text = slide_data['content']['main_concept']
        p.level = 0
        self._format_paragraph(p, theme, is_header=True)
        
        # Explanation
        p = text_frame.add_paragraph()
        p.text = slide_data['content']['explanation']
        p.level = 0
        self._format_paragraph(p, theme)
    
    def _create_activity_slide(self, prs: Presentation, slide_data: Dict[str, Any], theme: str):
        """Create activity slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = slide_data['title']
        self._format_title_text(title, theme)
        
        # Content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        
        content = slide_data['content']
        
        # Activity type and instructions
        p = text_frame.paragraphs[0]
        p.text = f"Activity: {content.get('activity_type', 'Exercise').replace('_', ' ').title()}"
        p.level = 0
        self._format_paragraph(p, theme, is_header=True)
        
        # Instructions
        p = text_frame.add_paragraph()
        p.text = content.get('instructions', 'Follow the instructions provided')
        p.level = 0
        self._format_paragraph(p, theme)
        
        # Materials if available
        if 'materials_needed' in content:
            p = text_frame.add_paragraph()
            p.text = "Materials needed:"
            p.level = 0
            self._format_paragraph(p, theme, is_header=True)
            
            for material in content['materials_needed']:
                p = text_frame.add_paragraph()
                p.text = material
                p.level = 1
                self._format_paragraph(p, theme)
    
    def _create_assessment_slide(self, prs: Presentation, slide_data: Dict[str, Any], theme: str):
        """Create assessment slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = slide_data['title']
        self._format_title_text(title, theme)
        
        # Content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        
        assessment = slide_data['content']
        
        # Add questions
        question_num = 1
        for question in assessment.get('questions', []):
            p = text_frame.paragraphs[0] if question_num == 1 else text_frame.add_paragraph()
            p.text = f"Q{question_num}: {question['question']}"
            p.level = 0
            self._format_paragraph(p, theme, is_header=True)
            
            # Add options for multiple choice
            if question['type'] == 'multiple_choice':
                for i, option in enumerate(question['options']):
                    p = text_frame.add_paragraph()
                    p.text = f"{chr(65+i)}. {option}"
                    p.level = 1
                    self._format_paragraph(p, theme)
            
            question_num += 1
            
            # Limit to 2 questions per slide
            if question_num > 2:
                break
    
    def _create_conclusion_slide(self, prs: Presentation, slide_data: Dict[str, Any], theme: str):
        """Create conclusion slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = slide_data['title']
        self._format_title_text(title, theme)
        
        # Content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        
        content = slide_data['content']
        
        # Key takeaways
        p = text_frame.paragraphs[0]
        p.text = "Key Takeaways:"
        p.level = 0
        self._format_paragraph(p, theme, is_header=True)
        
        for takeaway in content.get('key_takeaways', []):
            p = text_frame.add_paragraph()
            p.text = takeaway
            p.level = 1
            self._format_paragraph(p, theme)
        
        # Next steps
        p = text_frame.add_paragraph()
        p.text = "Next Steps:"
        p.level = 0
        self._format_paragraph(p, theme, is_header=True)
        
        for step in content.get('next_steps', []):
            p = text_frame.add_paragraph()
            p.text = step
            p.level = 1
            self._format_paragraph(p, theme)
    
    def _create_generic_slide(self, prs: Presentation, slide_data: Dict[str, Any], theme: str):
        """Create generic content slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = slide_data.get('title', 'Content')
        self._format_title_text(title, theme)
        
        # Content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        
        # Add content as text
        content_text = str(slide_data.get('content', ''))
        p = text_frame.paragraphs[0]
        p.text = content_text
        self._format_paragraph(p, theme)
    
    def _format_title_text(self, title_shape, theme: str):
        """Format title text"""
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
    
    def _format_subtitle_text(self, subtitle_shape, theme: str):
        """Format subtitle text"""
        for paragraph in subtitle_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = self.colors['text']
    
    def _format_paragraph(self, paragraph, theme: str, is_header: bool = False, is_intro: bool = False):
        """Format paragraph text"""
        if is_header:
            paragraph.font.size = Pt(20)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.colors['secondary']
        elif is_intro:
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.colors['text']
        else:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = self.colors['text']


class PDFGenerator:
    """Generate PDF documents from content data"""
    
    def __init__(self):
        setup_matplotlib_for_plotting()
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()
    
    def _setup_custom_styles(self):
        """Setup custom paragraph styles"""
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Title'],
            fontSize=24,
            spaceAfter=20,
            textColor=colors.HexColor('#365F92')
        ))
        
        self.styles.add(ParagraphStyle(
            name='CustomHeading',
            parent=self.styles['Heading1'],
            fontSize=18,
            spaceAfter=12,
            textColor=colors.HexColor('#4F81BD')
        ))
        
        self.styles.add(ParagraphStyle(
            name='CustomBody',
            parent=self.styles['Normal'],
            fontSize=12,
            spaceAfter=6,
            leftIndent=0
        ))
        
        self.styles.add(ParagraphStyle(
            name='CustomBullet',
            parent=self.styles['Normal'],
            fontSize=12,
            spaceAfter=6,
            leftIndent=20,
            bulletIndent=10
        ))
    
    def create_training_manual(self, 
                             module_data: Dict[str, Any], 
                             output_path: str,
                             include_charts: bool = True) -> str:
        """
        Create PDF training manual from module data
        
        Args:
            module_data: Training module data
            output_path: Path to save PDF
            include_charts: Whether to include charts and visualizations
            
        Returns:
            Path to created PDF
        """
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # Create PDF document
        doc = SimpleDocTemplate(
            output_path,
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=18
        )
        
        # Build content
        content = []
        
        # Title page
        content.extend(self._create_title_page(module_data))
        content.append(PageBreak())
        
        # Table of contents
        content.extend(self._create_table_of_contents(module_data))
        content.append(PageBreak())
        
        # Learning objectives
        content.extend(self._create_objectives_section(module_data))
        content.append(Spacer(1, 20))
        
        # Content sections
        content.extend(self._create_content_sections(module_data))
        
        # Activities
        if module_data.get('activities'):
            content.append(PageBreak())
            content.extend(self._create_activities_section(module_data))
        
        # Assessment
        if module_data.get('assessment'):
            content.append(PageBreak())
            content.extend(self._create_assessment_section(module_data))
        
        # Resources
        content.append(PageBreak())
        content.extend(self._create_resources_section(module_data))
        
        # Build PDF
        doc.build(content)
        
        return output_path
    
    def _create_title_page(self, module_data: Dict[str, Any]) -> List:
        """Create title page content"""
        content = []
        
        # Main title
        title = Paragraph(module_data['metadata']['title'], self.styles['CustomTitle'])
        content.append(title)
        content.append(Spacer(1, 30))
        
        # Metadata table
        metadata_data = [
            ['Duration:', f"{module_data['metadata']['duration_minutes']} minutes"],
            ['Learning Level:', module_data['metadata']['learning_level'].title()],
            ['Format:', module_data['metadata']['format_type'].title()],
            ['Created:', module_data['metadata']['created_date'][:10]]
        ]
        
        metadata_table = Table(metadata_data, colWidths=[2*inch, 3*inch])
        metadata_table.setStyle(TableStyle([
            ('FONTSIZE', (0, 0), (-1, -1), 12),
            ('ALIGN', (0, 0), (0, -1), 'RIGHT'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('GRID', (0, 0), (-1, -1), 1, colors.lightgrey)
        ]))
        
        content.append(metadata_table)
        content.append(Spacer(1, 50))
        
        # Description
        description = Paragraph(
            "This training manual provides comprehensive coverage of the topic with practical examples, "
            "interactive activities, and assessment opportunities.",
            self.styles['CustomBody']
        )
        content.append(description)
        
        return content
    
    def _create_table_of_contents(self, module_data: Dict[str, Any]) -> List:
        """Create table of contents"""
        content = []
        
        content.append(Paragraph("Table of Contents", self.styles['CustomTitle']))
        content.append(Spacer(1, 20))
        
        toc_items = [
            "1. Learning Objectives",
            "2. Content Overview",
            "3. Detailed Content"
        ]
        
        if module_data.get('activities'):
            toc_items.append("4. Activities")
        
        if module_data.get('assessment'):
            toc_items.append(f"{len(toc_items)+1}. Assessment")
        
        toc_items.append(f"{len(toc_items)+1}. Resources")
        
        for item in toc_items:
            content.append(Paragraph(item, self.styles['CustomBody']))
            content.append(Spacer(1, 6))
        
        return content
    
    def _create_objectives_section(self, module_data: Dict[str, Any]) -> List:
        """Create learning objectives section"""
        content = []
        
        content.append(Paragraph("Learning Objectives", self.styles['CustomTitle']))
        content.append(Spacer(1, 15))
        
        content.append(Paragraph(
            "By the end of this training module, you will be able to:",
            self.styles['CustomBody']
        ))
        content.append(Spacer(1, 10))
        
        for objective in module_data['learning_objectives']:
            bullet_text = f"• {objective}"
            content.append(Paragraph(bullet_text, self.styles['CustomBullet']))
        
        return content
    
    def _create_content_sections(self, module_data: Dict[str, Any]) -> List:
        """Create detailed content sections"""
        content = []
        
        content.append(Paragraph("Training Content", self.styles['CustomTitle']))
        content.append(Spacer(1, 15))
        
        # Content outline
        content.append(Paragraph("Content Overview", self.styles['CustomHeading']))
        
        for section in module_data['content_outline']:
            section_title = f"{section['section']} ({section['duration_minutes']} minutes)"
            content.append(Paragraph(section_title, self.styles['CustomBody']))
            
            for point in section['key_points']:
                bullet_text = f"• {point}"
                content.append(Paragraph(bullet_text, self.styles['CustomBullet']))
            
            content.append(Spacer(1, 10))
        
        # Detailed content
        content.append(Spacer(1, 20))
        content.append(Paragraph("Detailed Content", self.styles['CustomHeading']))
        
        for section in module_data['detailed_content']:
            if section['type'] != 'activity':  # Activities handled separately
                content.append(Paragraph(section['title'], self.styles['CustomHeading']))
                
                if 'content' in section:
                    section_content = section['content']
                    
                    # Add explanation
                    if 'explanation' in section_content:
                        content.append(Paragraph(section_content['explanation'], self.styles['CustomBody']))
                        content.append(Spacer(1, 10))
                    
                    # Add examples
                    if 'examples' in section_content:
                        content.append(Paragraph("Examples:", self.styles['CustomBody']))
                        for example in section_content['examples']:
                            bullet_text = f"• {example}"
                            content.append(Paragraph(bullet_text, self.styles['CustomBullet']))
                        content.append(Spacer(1, 10))
                    
                    # Add best practices
                    if 'best_practices' in section_content:
                        content.append(Paragraph("Best Practices:", self.styles['CustomBody']))
                        for practice in section_content['best_practices']:
                            bullet_text = f"• {practice}"
                            content.append(Paragraph(bullet_text, self.styles['CustomBullet']))
                
                content.append(Spacer(1, 15))
        
        return content
    
    def _create_activities_section(self, module_data: Dict[str, Any]) -> List:
        """Create activities section"""
        content = []
        
        content.append(Paragraph("Interactive Activities", self.styles['CustomTitle']))
        content.append(Spacer(1, 15))
        
        for activity in module_data['activities']:
            content.append(Paragraph(activity['title'], self.styles['CustomHeading']))
            
            # Activity details
            details = [
                ['Type:', activity['type'].replace('_', ' ').title()],
                ['Duration:', f"{activity['duration_minutes']} minutes"],
                ['Description:', activity['description']]
            ]
            
            details_table = Table(details, colWidths=[1.5*inch, 4.5*inch])
            details_table.setStyle(TableStyle([
                ('FONTSIZE', (0, 0), (-1, -1), 10),
                ('ALIGN', (0, 0), (0, -1), 'RIGHT'),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('GRID', (0, 0), (-1, -1), 1, colors.lightgrey)
            ]))
            
            content.append(details_table)
            content.append(Spacer(1, 10))
            
            # Instructions
            content.append(Paragraph("Instructions:", self.styles['CustomBody']))
            for instruction in activity['instructions']:
                bullet_text = f"• {instruction}"
                content.append(Paragraph(bullet_text, self.styles['CustomBullet']))
            
            content.append(Spacer(1, 20))
        
        return content
    
    def _create_assessment_section(self, module_data: Dict[str, Any]) -> List:
        """Create assessment section"""
        content = []
        
        content.append(Paragraph("Assessment", self.styles['CustomTitle']))
        content.append(Spacer(1, 15))
        
        assessment = module_data['assessment']
        
        content.append(Paragraph(
            "Test your understanding with the following questions:",
            self.styles['CustomBody']
        ))
        content.append(Spacer(1, 15))
        
        for i, question in enumerate(assessment['questions'], 1):
            content.append(Paragraph(f"Question {i}:", self.styles['CustomHeading']))
            content.append(Paragraph(question['question'], self.styles['CustomBody']))
            
            if question['type'] == 'multiple_choice':
                content.append(Spacer(1, 10))
                for j, option in enumerate(question['options']):
                    option_text = f"{chr(65+j)}. {option}"
                    content.append(Paragraph(option_text, self.styles['CustomBullet']))
            
            content.append(Spacer(1, 15))
        
        return content
    
    def _create_resources_section(self, module_data: Dict[str, Any]) -> List:
        """Create resources section"""
        content = []
        
        content.append(Paragraph("Additional Resources", self.styles['CustomTitle']))
        content.append(Spacer(1, 15))
        
        if module_data.get('resources'):
            content.append(Paragraph(
                "The following resources were used in creating this training module:",
                self.styles['CustomBody']
            ))
            content.append(Spacer(1, 10))
            
            for resource in module_data['resources']:
                resource_text = f"• {resource['title']} ({resource['type']})"
                content.append(Paragraph(resource_text, self.styles['CustomBullet']))
                
                if 'description' in resource:
                    description_text = f"  {resource['description']}"
                    content.append(Paragraph(description_text, self.styles['CustomBody']))
                
                content.append(Spacer(1, 5))
        else:
            content.append(Paragraph(
                "No additional resources were identified for this module.",
                self.styles['CustomBody']
            ))
        
        return content